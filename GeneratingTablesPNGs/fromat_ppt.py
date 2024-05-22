from pptx import Presentation
from pptx.util import Cm

# Load the presentation
presentation_path = 'Task2_format.pptx'
presentation = Presentation(presentation_path)

# Set the desired picture formatting
height = Cm(17.67)
width = Cm(33.87)
rotation = 0
scale_height = 0.73
scale_width = 0.77
horizontal_position = Cm(0)
vertical_position = Cm(0.7)

# Function to format picture
def format_picture(picture):
    picture.height = height
    picture.width = width
    picture.rotation = rotation
    picture.left = horizontal_position
    picture.top = vertical_position

# Loop through slides and shapes to find and format pictures
for slide in presentation.slides:
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture shape type is 13
            format_picture(shape)

# Save the modified presentation
presentation.save('formatted_task2.pptx')
